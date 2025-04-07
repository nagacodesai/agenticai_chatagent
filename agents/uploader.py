import os
import sys
import argparse
import pandas as pd
import requests
from tqdm import tqdm
from pinecone import Pinecone, ServerlessSpec

# ➕ Add parent directory to path so we can import sibling modules
sys.path.append(os.path.abspath(os.path.join(os.path.dirname(__file__), '..')))

# 🔐 Load secure Pinecone credentials from env_loader
from agents.env_loader import (
    PINECONE_API_KEY,
    PINECONE_ENV,
    PINECONE_INDEX
)

from agents.embedder import get_embedding

# 🔧 Initialize Pinecone
pc = Pinecone(api_key=PINECONE_API_KEY)

# 📌 Create index if it doesn't exist
if PINECONE_INDEX not in pc.list_indexes().names():
    pc.create_index(
        name=PINECONE_INDEX,
        dimension=1536,
        metric="cosine",
        spec=ServerlessSpec(cloud="aws", region=PINECONE_ENV)
    )
index = pc.Index(PINECONE_INDEX)

# =====================
# 📥 Upload Entry Points
# =====================

def upload_from_csv(file_path: str):
    """Upload tariff data from a local CSV file."""
    df = pd.read_csv(file_path)
    df = _clean_tariff_data(df)
    _process_dataframe(df)

def upload_from_excel(file_path: str):
    """Upload tariff data from a local Excel file."""
    df = pd.read_excel(file_path)
    df = _clean_tariff_data(df)
    _process_dataframe(df)

def upload_from_pdf(file_path: str):
    """Extract text from PDF file and upload as embedded documents."""
    from PyPDF2 import PdfReader
    reader = PdfReader(file_path)
    texts = [page.extract_text() for page in reader.pages if page.extract_text()]
    _embed_text_chunks(texts)

def upload_from_word(file_path: str):
    """Extract text from Word (.docx) file and upload as embedded documents."""
    from docx import Document
    doc = Document(file_path)
    texts = [para.text for para in doc.paragraphs if para.text.strip()]
    _embed_text_chunks(texts)

def upload_from_ppt(file_path: str):
    """Extract text from PowerPoint (.pptx) file and upload as embedded documents."""
    from pptx import Presentation
    prs = Presentation(file_path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
    _embed_text_chunks(texts)

def upload_from_api(api_url: str):
    """Fetch JSON data from API endpoint and upload tariff info."""
    response = requests.get(api_url)
    response.raise_for_status()
    data = response.json()
    df = pd.DataFrame(data)
    df = _clean_tariff_data(df)
    _process_dataframe(df)

def upload_from_dataframe(df: pd.DataFrame):
    """Directly upload from a DataFrame — useful for pre-processed dynamic data."""
    print("📄 Uploading from provided DataFrame...")
    df = _clean_tariff_data(df)
    _process_dataframe(df)

# =====================
# 🔧 Helper Functions
# =====================

def _clean_tariff_data(df: pd.DataFrame):
    """Remove percentage signs and convert tariff columns to floats."""
    if "TariffsCharged2USA" in df.columns:
        df["TariffsCharged2USA"] = (
            df["TariffsCharged2USA"].astype(str).str.replace('%', '').astype(float)
        )
    if "USAReciprocalTariffs" in df.columns:
        df["USAReciprocalTariffs"] = (
            df["USAReciprocalTariffs"].astype(str).str.replace('%', '').astype(float)
        )
    return df

def _process_dataframe(df: pd.DataFrame):
    """Convert each row of DataFrame into a Pinecone vector with metadata."""
    vectors = []
    for i, row in tqdm(df.iterrows(), total=len(df)):
        doc_text = ", ".join(f"{col}: {row[col]}" for col in df.columns)
        vector = get_embedding(doc_text)
        vectors.append({
            "id": f"row-{i}",
            "values": vector,
            "metadata": {"text": doc_text}
        })
    _upload_vectors(vectors)

def _embed_text_chunks(texts: list):
    """Embed plain text chunks (from PDF, Word, PPT, etc.) and upload."""
    vectors = []
    for i, text in tqdm(enumerate(texts), total=len(texts)):
        vector = get_embedding(text)
        vectors.append({
            "id": f"doc-{i}",
            "values": vector,
            "metadata": {"text": text}
        })
    _upload_vectors(vectors)

def _upload_vectors(vectors: list):
    """Upsert vector chunks to Pinecone in batches."""
    from more_itertools import chunked
    for chunk in chunked(vectors, 100):
        index.upsert(vectors=chunk)
    print("✅ Successfully upserted into Pinecone.")

# =====================
# 🎛️ CLI Interface
# =====================

def main():
    """Command-line interface to support uploads from different formats."""
    parser = argparse.ArgumentParser(description="Upload embeddings to Pinecone from various sources.")
    group = parser.add_mutually_exclusive_group(required=True)
    group.add_argument("--csv", help="Path to CSV file")
    group.add_argument("--excel", help="Path to Excel file")
    group.add_argument("--pdf", help="Path to PDF file")
    group.add_argument("--word", help="Path to Word file (.docx)")
    group.add_argument("--ppt", help="Path to PowerPoint file (.pptx)")
    group.add_argument("--api", help="GET API URL returning JSON")

    args = parser.parse_args()

    if args.csv:
        upload_from_csv(args.csv)
    elif args.excel:
        upload_from_excel(args.excel)
    elif args.pdf:
        upload_from_pdf(args.pdf)
    elif args.word:
        upload_from_word(args.word)
    elif args.ppt:
        upload_from_ppt(args.ppt)
    elif args.api:
        upload_from_api(args.api)

# 🏁 Run when executed directly
if __name__ == "__main__":
    main()
